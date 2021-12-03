#the purpose of this file is to make powerpoint presentations out of lab data
#sheets without requiring any transposition

import os
import re
from pptx import Presentation as pres
from pptx.util import Inches
import csv
import datetime
import shutil

class powerPoint:
    def __init__(self):
        self.title = "Blank place holder"
        self.prs = pres()
        self.picture = self.prs.slide_layouts[8]
        self.title = self.prs.slide_layouts[0]
        self.title_content = self.prs.slide_layouts[1]

        
        
    def make_title(self):
        slide = self.prs.slides.add_slide(self.title)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = ("Summary of Data "
                      + datetime.datetime.today().strftime('%Y-%m-%d'))
        subtitle.text = "Nick Gregory"

    def get_images(self):
        d = datetime.datetime.today().strftime('%Y-%m-%d')
        cwd = os.getcwd()
        files = os.listdir(cwd)
        txtRe = re.compile('(?P<file>.*).png', re.I)
        #print (txtRe)
        for file in files:
            match = txtRe.match(file)
            #print (file + ": ")
            #print (match)
            if match:
                slide = self.prs.slides.add_slide(self.picture)
                title = slide.shapes.title
                file_name_parts = file.split('-')[0].split('_') 
                filename = " ".join(file_name_parts)

                title.text = self.get_graph_title(file)
                placeholder  = slide.placeholders[1]
                picture = placeholder.insert_picture(file)
                #no cropping
                picture.crop_top = 0
                picture.crop_left = 0
                picture.crop_bottom = 0
                picture.crop_right = 0
                if os.path.isdir(cwd+"\\"+d):
                    shutil.move(file, cwd+"\\"+d+"\\")
                else:
                    os.mkdir(cwd+"\\"+d)
                    shutil.move(file, cwd+"\\"+d+"\\")
                
    def test_layouts(self):
        x = 8
        while x <= 8:
            print(x)
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[x])
            title = slide.shapes.title
            subtitle = slide.placeholders[2]
            title.text = ("Lab Meeting " 
                          + datetime.datetime.today().strftime('%Y-%m-%d'))
            subtitle.text = "Nick Gregory"
            placeholder  = slide.placeholders[1]
            picture = placeholder.insert_picture("PR_boxplot_2019-01-03.png")
            x+=1

    def get_graph_title(self, filename):
        txtRe = re.compile('\w* (.*) by (.*).png', re.I)
        match = txtRe.match(filename)
        graph_title = f'Plot: {match[1]} by {match[2]}'
        return graph_title






if __name__ == "__main__":
    cwd = os.getcwd()
    d = datetime.datetime.today().strftime('%Y-%m-%d')
    presentation = powerPoint()
    presentation.make_title()
    presentation.get_images()  # Moves the "gotten" images to the folder as it loads
    # them into the pptx.
    parameter = 'bdi sum'
    name = f"{parameter} summary.pptx"
    presentation.prs.save(name)
    shutil.move(name, cwd + "\\" + d + "\\")